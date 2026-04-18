"""
Open Brain File Ingestion — Content Extraction Service

FastAPI service that extracts text and metadata from documents.
Supports: PDF, DOCX, PPTX, XLSX, TXT, MD, CSV, HTML.

Memory budget: 1.5GB RSS ceiling. All extraction uses streaming/chunked
reads where possible. Large files are processed page-by-page or sheet-by-sheet.
"""

import csv
import logging
import mimetypes
import os
import time
from pathlib import Path
from typing import Any

import xxhash
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=os.getenv("LOG_LEVEL", "INFO").upper(),
    format="%(asctime)s %(levelname)s %(name)s: %(message)s",
)
logger = logging.getLogger("file-ingestion")

# ---------------------------------------------------------------------------
# FastAPI app
# ---------------------------------------------------------------------------
app = FastAPI(
    title="Open Brain File Ingestion",
    version="1.0.0",
    description="Content extraction service for Office/document formats",
)

# ---------------------------------------------------------------------------
# Models
# ---------------------------------------------------------------------------


class ExtractionRequest(BaseModel):
    file_path: str = Field(..., description="Absolute path to the file to extract")


class SectionInfo(BaseModel):
    title: str
    content: str
    level: int = 1


class ExtractionResult(BaseModel):
    text: str = Field(..., description="Extracted plain text content")
    metadata: dict[str, Any] = Field(default_factory=dict, description="Document metadata")
    sections: list[SectionInfo] = Field(
        default_factory=list, description="Document sections (if applicable)"
    )
    content_hash: str = Field(..., description="xxhash of extracted text")
    file_size: int = Field(..., description="Original file size in bytes")
    mime_type: str = Field(..., description="Detected MIME type")
    extraction_time_ms: int = Field(..., description="Extraction duration in milliseconds")


class HealthResponse(BaseModel):
    status: str
    version: str
    supported_types: list[str]


# ---------------------------------------------------------------------------
# Supported extensions → MIME type mapping
# ---------------------------------------------------------------------------
SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".html": "text/html",
    ".htm": "text/html",
}

# Maximum text size to prevent memory blowout (50MB of text is ~50M chars)
MAX_TEXT_SIZE = 50 * 1024 * 1024


# ---------------------------------------------------------------------------
# Extractors — one per format
# ---------------------------------------------------------------------------


def _extract_pdf(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from PDF using pdfplumber. Streams page-by-page."""
    import pdfplumber

    text_parts: list[str] = []
    metadata: dict[str, Any] = {}
    sections: list[SectionInfo] = []
    total_chars = 0

    with pdfplumber.open(file_path) as pdf:
        metadata["page_count"] = len(pdf.pages)
        if pdf.metadata:
            for key in ("Title", "Author", "Subject", "Creator", "Producer"):
                val = pdf.metadata.get(key)
                if val:
                    metadata[key.lower()] = val

        for i, page in enumerate(pdf.pages):
            if total_chars > MAX_TEXT_SIZE:
                metadata["truncated"] = True
                break
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
            total_chars += len(page_text)

            if page_text.strip():
                sections.append(
                    SectionInfo(
                        title=f"Page {i + 1}",
                        content=page_text,
                        level=1,
                    )
                )

    return "\n\n".join(text_parts), metadata, sections


def _extract_docx(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(str(file_path))
    metadata: dict[str, Any] = {}
    sections: list[SectionInfo] = []
    text_parts: list[str] = []

    # Core properties
    props = doc.core_properties
    if props.title:
        metadata["title"] = props.title
    if props.author:
        metadata["author"] = props.author
    if props.subject:
        metadata["subject"] = props.subject
    if props.modified:
        metadata["modified"] = str(props.modified)

    current_section_title = ""
    current_section_parts: list[str] = []
    current_level = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        text_parts.append(text)

        # Track heading-based sections
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            # Flush previous section
            if current_section_parts:
                sections.append(
                    SectionInfo(
                        title=current_section_title or "Introduction",
                        content="\n".join(current_section_parts),
                        level=current_level,
                    )
                )
                current_section_parts = []

            current_section_title = text
            try:
                current_level = int(para.style.name.replace("Heading ", "").replace("Heading", "1"))
            except (ValueError, AttributeError):
                current_level = 1
        else:
            current_section_parts.append(text)

    # Flush last section
    if current_section_parts:
        sections.append(
            SectionInfo(
                title=current_section_title or "Content",
                content="\n".join(current_section_parts),
                level=current_level,
            )
        )

    metadata["paragraph_count"] = len(doc.paragraphs)

    return "\n\n".join(text_parts), metadata, sections


def _extract_pptx(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    metadata: dict[str, Any] = {}
    sections: list[SectionInfo] = []
    text_parts: list[str] = []

    props = prs.core_properties
    if props.title:
        metadata["title"] = props.title
    if props.author:
        metadata["author"] = props.author

    metadata["slide_count"] = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Also extract from tables
            if shape.has_table:
                for row in shape.table.rows:  # type: ignore[attr-defined]
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))

        slide_text = "\n".join(slide_texts)
        if slide_text:
            text_parts.append(slide_text)
            sections.append(
                SectionInfo(
                    title=f"Slide {i + 1}",
                    content=slide_text,
                    level=1,
                )
            )

    return "\n\n---\n\n".join(text_parts), metadata, sections


def _extract_xlsx(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from XLSX using openpyxl. Read-only mode for memory efficiency."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    metadata: dict[str, Any] = {}
    sections: list[SectionInfo] = []
    text_parts: list[str] = []

    metadata["sheet_count"] = len(wb.sheetnames)
    metadata["sheet_names"] = wb.sheetnames

    if wb.properties:
        if wb.properties.title:
            metadata["title"] = wb.properties.title
        if wb.properties.creator:
            metadata["author"] = wb.properties.creator

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_rows: list[str] = []

        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                sheet_rows.append(" | ".join(cells))

        sheet_text = "\n".join(sheet_rows)
        if sheet_text:
            text_parts.append(f"[Sheet: {sheet_name}]\n{sheet_text}")
            sections.append(
                SectionInfo(
                    title=f"Sheet: {sheet_name}",
                    content=sheet_text,
                    level=1,
                )
            )

    wb.close()
    return "\n\n".join(text_parts), metadata, sections


def _extract_text(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from plain text files (TXT, MD)."""
    metadata: dict[str, Any] = {}

    with open(file_path, encoding="utf-8", errors="replace") as f:
        text = f.read(MAX_TEXT_SIZE)

    if len(text) == MAX_TEXT_SIZE:
        metadata["truncated"] = True

    line_count = text.count("\n") + 1
    metadata["line_count"] = line_count

    # For markdown, extract heading-based sections
    sections: list[SectionInfo] = []
    ext = file_path.suffix.lower()
    if ext == ".md":
        current_title = ""
        current_parts: list[str] = []
        current_level = 1

        for line in text.split("\n"):
            if line.startswith("#"):
                if current_parts:
                    sections.append(
                        SectionInfo(
                            title=current_title or "Introduction",
                            content="\n".join(current_parts),
                            level=current_level,
                        )
                    )
                    current_parts = []
                hashes = len(line) - len(line.lstrip("#"))
                current_title = line.lstrip("#").strip()
                current_level = min(hashes, 6)
            else:
                current_parts.append(line)

        if current_parts:
            sections.append(
                SectionInfo(
                    title=current_title or "Content",
                    content="\n".join(current_parts),
                    level=current_level,
                )
            )

    return text, metadata, sections


def _extract_csv(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from CSV files."""
    metadata: dict[str, Any] = {}
    rows: list[str] = []

    with open(file_path, encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row_count, row in enumerate(reader):
            if row_count > 100_000:
                metadata["truncated"] = True
                break
            rows.append(" | ".join(row))

    metadata["row_count"] = len(rows)

    # First row is typically header
    header = rows[0] if rows else ""
    metadata["columns"] = header

    return "\n".join(rows), metadata, []


def _extract_html(file_path: Path) -> tuple[str, dict[str, Any], list[SectionInfo]]:
    """Extract text from HTML files using BeautifulSoup."""
    from bs4 import BeautifulSoup

    metadata: dict[str, Any] = {}
    sections: list[SectionInfo] = []

    with open(file_path, encoding="utf-8", errors="replace") as f:
        content = f.read(MAX_TEXT_SIZE)

    if len(content) == MAX_TEXT_SIZE:
        metadata["truncated"] = True

    soup = BeautifulSoup(content, "html.parser")

    # Extract title
    title_tag = soup.find("title")
    if title_tag and title_tag.string:
        metadata["title"] = title_tag.string.strip()

    # Extract meta tags (bs4 .get() may return str | list[str] | None)
    for meta in soup.find_all("meta"):
        name_raw = meta.get("name", "") if hasattr(meta, "get") else ""  # type: ignore[union-attr]
        content_raw = meta.get("content", "") if hasattr(meta, "get") else ""  # type: ignore[union-attr]
        name = (name_raw if isinstance(name_raw, str) else "").lower()
        content_attr = content_raw if isinstance(content_raw, str) else ""
        if name in ("author", "description", "keywords") and content_attr:
            metadata[name] = content_attr

    # Extract text
    text = soup.get_text(separator="\n", strip=True)

    # Extract heading-based sections
    for heading in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
        level = int(heading.name[1])
        heading_text = heading.get_text(strip=True)
        # Collect text until next heading
        content_parts: list[str] = []
        for sibling in heading.find_next_siblings():
            if sibling.name and sibling.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
                break
            sib_text = sibling.get_text(strip=True)
            if sib_text:
                content_parts.append(sib_text)
        if heading_text:
            sections.append(
                SectionInfo(
                    title=heading_text,
                    content="\n".join(content_parts),
                    level=level,
                )
            )

    return text, metadata, sections


# Extractor dispatch
EXTRACTORS: dict[str, Any] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".txt": _extract_text,
    ".md": _extract_text,
    ".csv": _extract_csv,
    ".html": _extract_html,
    ".htm": _extract_html,
}


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------


@app.get("/health", response_model=HealthResponse)
async def health():
    """Health check endpoint."""
    return HealthResponse(
        status="healthy",
        version="1.0.0",
        supported_types=list(SUPPORTED_EXTENSIONS.keys()),
    )


@app.post("/extract", response_model=ExtractionResult)
async def extract(req: ExtractionRequest):
    """Extract text and metadata from a document file.

    Accepts a file path (must be accessible to the container via mounted volume)
    and returns extracted text, metadata, and structural sections.
    """
    file_path = Path(req.file_path)

    # Validate file exists
    if not file_path.exists():
        raise HTTPException(status_code=404, detail=f"File not found: {req.file_path}")

    if not file_path.is_file():
        raise HTTPException(status_code=400, detail=f"Not a file: {req.file_path}")

    # Determine file type
    ext = file_path.suffix.lower()
    if ext not in EXTRACTORS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Supported: {list(SUPPORTED_EXTENSIONS.keys())}",
        )

    # Get file size
    file_size = file_path.stat().st_size

    # Detect MIME type
    mime_type = SUPPORTED_EXTENSIONS.get(ext, "")
    if not mime_type:
        mime_type, _ = mimetypes.guess_type(str(file_path))
        mime_type = mime_type or "application/octet-stream"

    # Extract content
    start = time.monotonic()
    try:
        extractor = EXTRACTORS[ext]
        text, metadata, sections = extractor(file_path)
    except Exception as e:
        logger.error("Extraction failed for %s: %s", file_path, e, exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Extraction failed: {type(e).__name__}: {e}",
        ) from e

    elapsed_ms = int((time.monotonic() - start) * 1000)

    # Compute content hash
    content_hash = xxhash.xxh64(text.encode("utf-8")).hexdigest()

    logger.info(
        "Extracted %s: %d chars, %d sections, %dms",
        file_path.name,
        len(text),
        len(sections),
        elapsed_ms,
    )

    return ExtractionResult(
        text=text,
        metadata=metadata,
        sections=sections,
        content_hash=content_hash,
        file_size=file_size,
        mime_type=mime_type,
        extraction_time_ms=elapsed_ms,
    )


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import uvicorn

    port = int(os.getenv("PORT", "8080"))
    uvicorn.run(
        "src.extract:app",
        host="0.0.0.0",
        port=port,
        log_level="info",
        workers=1,  # Single worker to stay within 1.5GB memory ceiling
    )
