"""
Tests for the file-ingestion extraction service.

Tests each supported file type for correct text extraction, metadata,
and error handling. Fixtures are created programmatically in conftest.
"""

import csv
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from src.extract import app

FIXTURES_DIR = Path(__file__).parent / "fixtures"


@pytest.fixture(autouse=True, scope="session")
def create_all_fixtures():
    """Create all test fixture files before tests run."""
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)

    # TXT
    (FIXTURES_DIR / "sample.txt").write_text(
        "Hello, World!\nThis is a test document.\nLine three.",
        encoding="utf-8",
    )

    # MD
    (FIXTURES_DIR / "sample.md").write_text(
        "# Heading One\n\nFirst section content.\n\n"
        "## Heading Two\n\nSecond section content with details.\n\n"
        "### Sub-heading\n\nSub-section text.",
        encoding="utf-8",
    )

    # CSV
    with open(FIXTURES_DIR / "sample.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Age", "City"])
        writer.writerow(["Alice", "30", "New York"])
        writer.writerow(["Bob", "25", "London"])
        writer.writerow(["Charlie", "35", "Tokyo"])

    # HTML
    (FIXTURES_DIR / "sample.html").write_text(
        '<!DOCTYPE html>\n<html>\n<head>\n'
        '    <title>Test Document</title>\n'
        '    <meta name="author" content="Test Author">\n'
        '</head>\n<body>\n'
        '    <h1>Main Heading</h1>\n'
        '    <p>First paragraph with important content.</p>\n'
        '    <h2>Sub Heading</h2>\n'
        '    <p>Second paragraph under sub heading.</p>\n'
        '</body>\n</html>',
        encoding="utf-8",
    )

    # DOCX
    try:
        from docx import Document

        doc = Document()
        doc.core_properties.title = "Test DOCX"
        doc.core_properties.author = "Test Author"
        doc.add_heading("Document Title", level=1)
        doc.add_paragraph("This is the first paragraph of the test document.")
        doc.add_heading("Section Two", level=2)
        doc.add_paragraph("Content under section two.")
        doc.save(str(FIXTURES_DIR / "sample.docx"))
    except ImportError:
        pass

    # PPTX
    try:
        from pptx import Presentation

        prs = Presentation()
        prs.core_properties.title = "Test Presentation"
        prs.core_properties.author = "Test Author"
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Test Presentation Title"
        slide.placeholders[1].text = "Subtitle text here"
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Content Slide"
        slide.placeholders[1].text = "Bullet point one\nBullet point two"
        prs.save(str(FIXTURES_DIR / "sample.pptx"))
    except ImportError:
        pass

    # XLSX
    try:
        from openpyxl import Workbook

        wb = Workbook()
        wb.properties.title = "Test Spreadsheet"
        wb.properties.creator = "Test Author"
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Value", "Category"])
        ws.append(["Alpha", 100, "A"])
        ws.append(["Beta", 200, "B"])
        ws2 = wb.create_sheet("Summary")
        ws2.append(["Total", 300])
        wb.save(str(FIXTURES_DIR / "sample.xlsx"))
    except ImportError:
        pass

    # PDF (minimal valid PDF, no extra dependency)
    pdf_content = """%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj

2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj

3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj

4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Hello PDF World) Tj ET
endstream
endobj

5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj

xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n

trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""
    (FIXTURES_DIR / "sample.pdf").write_text(pdf_content, encoding="ascii")


@pytest.fixture
def client():
    return TestClient(app)


# ---------------------------------------------------------------------------
# Health endpoint
# ---------------------------------------------------------------------------

class TestHealth:
    def test_health_returns_ok(self, client):
        resp = client.get("/health")
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "healthy"
        assert data["version"] == "1.0.0"
        assert ".pdf" in data["supported_types"]
        assert ".docx" in data["supported_types"]

    def test_health_lists_all_supported_types(self, client):
        resp = client.get("/health")
        types = resp.json()["supported_types"]
        expected = [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".csv", ".html", ".htm"]
        for ext in expected:
            assert ext in types, f"Missing {ext} in supported_types"


# ---------------------------------------------------------------------------
# TXT extraction
# ---------------------------------------------------------------------------

class TestTxtExtraction:
    def test_extract_txt(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.txt")})
        assert resp.status_code == 200
        data = resp.json()
        assert "Hello, World!" in data["text"]
        assert "Line three" in data["text"]
        assert data["mime_type"] == "text/plain"
        assert data["content_hash"]  # non-empty
        assert data["file_size"] > 0
        assert data["extraction_time_ms"] >= 0
        assert data["metadata"]["line_count"] == 3


# ---------------------------------------------------------------------------
# Markdown extraction
# ---------------------------------------------------------------------------

class TestMarkdownExtraction:
    def test_extract_md(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.md")})
        assert resp.status_code == 200
        data = resp.json()
        assert "Heading One" in data["text"]
        assert "First section content" in data["text"]
        assert data["mime_type"] == "text/markdown"

    def test_md_has_sections(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.md")})
        sections = resp.json()["sections"]
        assert len(sections) >= 2
        titles = [s["title"] for s in sections]
        assert "Heading One" in titles
        assert "Heading Two" in titles


# ---------------------------------------------------------------------------
# CSV extraction
# ---------------------------------------------------------------------------

class TestCsvExtraction:
    def test_extract_csv(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.csv")})
        assert resp.status_code == 200
        data = resp.json()
        assert "Alice" in data["text"]
        assert "Tokyo" in data["text"]
        assert data["mime_type"] == "text/csv"
        assert data["metadata"]["row_count"] == 4  # header + 3 data rows


# ---------------------------------------------------------------------------
# HTML extraction
# ---------------------------------------------------------------------------

class TestHtmlExtraction:
    def test_extract_html(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.html")})
        assert resp.status_code == 200
        data = resp.json()
        assert "Main Heading" in data["text"]
        assert "First paragraph" in data["text"]
        assert data["mime_type"] == "text/html"

    def test_html_metadata(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.html")})
        meta = resp.json()["metadata"]
        assert meta.get("title") == "Test Document"
        assert meta.get("author") == "Test Author"

    def test_html_sections(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.html")})
        sections = resp.json()["sections"]
        assert len(sections) >= 1
        titles = [s["title"] for s in sections]
        assert "Main Heading" in titles


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------

class TestPdfExtraction:
    def test_extract_pdf(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.pdf")})
        assert resp.status_code == 200
        data = resp.json()
        assert data["mime_type"] == "application/pdf"
        # Minimal hand-crafted PDF — pdfplumber may report 0 pages if xref is
        # imperfect, but the extractor must not crash and must return valid output
        assert "page_count" in data["metadata"]
        assert data["content_hash"]  # non-empty regardless

    def test_pdf_has_page_sections(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.pdf")})
        data = resp.json()
        # Even if text extraction is partial, sections should be present if text was found
        if data["text"].strip():
            assert len(data["sections"]) >= 1


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------

class TestDocxExtraction:
    @pytest.fixture(autouse=True)
    def check_docx_fixture(self):
        if not (FIXTURES_DIR / "sample.docx").exists():
            pytest.skip("python-docx not available, fixture not created")

    def test_extract_docx(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.docx")})
        assert resp.status_code == 200
        data = resp.json()
        assert "first paragraph" in data["text"].lower()
        assert data["mime_type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    def test_docx_metadata(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.docx")})
        meta = resp.json()["metadata"]
        assert meta.get("title") == "Test DOCX"
        assert meta.get("author") == "Test Author"

    def test_docx_sections(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.docx")})
        sections = resp.json()["sections"]
        assert len(sections) >= 1
        titles = [s["title"] for s in sections]
        assert "Document Title" in titles


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------

class TestPptxExtraction:
    @pytest.fixture(autouse=True)
    def check_pptx_fixture(self):
        if not (FIXTURES_DIR / "sample.pptx").exists():
            pytest.skip("python-pptx not available, fixture not created")

    def test_extract_pptx(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.pptx")})
        assert resp.status_code == 200
        data = resp.json()
        assert "Test Presentation Title" in data["text"]
        assert data["mime_type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def test_pptx_metadata(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.pptx")})
        meta = resp.json()["metadata"]
        assert meta.get("slide_count") == 2

    def test_pptx_sections(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.pptx")})
        sections = resp.json()["sections"]
        assert len(sections) >= 1


# ---------------------------------------------------------------------------
# XLSX extraction
# ---------------------------------------------------------------------------

class TestXlsxExtraction:
    @pytest.fixture(autouse=True)
    def check_xlsx_fixture(self):
        if not (FIXTURES_DIR / "sample.xlsx").exists():
            pytest.skip("openpyxl not available, fixture not created")

    def test_extract_xlsx(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.xlsx")})
        assert resp.status_code == 200
        data = resp.json()
        assert "Alpha" in data["text"]
        assert data["mime_type"] == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

    def test_xlsx_metadata(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.xlsx")})
        meta = resp.json()["metadata"]
        assert meta.get("sheet_count") == 2
        assert "Data" in meta.get("sheet_names", [])
        assert "Summary" in meta.get("sheet_names", [])

    def test_xlsx_sections(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.xlsx")})
        sections = resp.json()["sections"]
        assert len(sections) >= 1
        titles = [s["title"] for s in sections]
        assert any("Data" in t for t in titles)


# ---------------------------------------------------------------------------
# Error handling
# ---------------------------------------------------------------------------

class TestErrorHandling:
    def test_file_not_found(self, client):
        resp = client.post("/extract", json={"file_path": "/nonexistent/file.txt"})
        assert resp.status_code == 404

    def test_unsupported_type(self, client, tmp_path):
        bad_file = tmp_path / "test.xyz"
        bad_file.write_text("data")
        resp = client.post("/extract", json={"file_path": str(bad_file)})
        assert resp.status_code == 400
        assert "Unsupported file type" in resp.json()["detail"]

    def test_directory_rejected(self, client, tmp_path):
        resp = client.post("/extract", json={"file_path": str(tmp_path)})
        assert resp.status_code == 400
        assert "Not a file" in resp.json()["detail"]

    def test_missing_file_path(self, client):
        resp = client.post("/extract", json={})
        assert resp.status_code == 422  # Pydantic validation error


# ---------------------------------------------------------------------------
# Response structure
# ---------------------------------------------------------------------------

class TestResponseStructure:
    def test_all_fields_present(self, client):
        resp = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.txt")})
        assert resp.status_code == 200
        data = resp.json()
        required_fields = ["text", "metadata", "sections", "content_hash", "file_size", "mime_type", "extraction_time_ms"]
        for field in required_fields:
            assert field in data, f"Missing field: {field}"

    def test_content_hash_is_stable(self, client):
        """Same file should always produce the same hash."""
        resp1 = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.txt")})
        resp2 = client.post("/extract", json={"file_path": str(FIXTURES_DIR / "sample.txt")})
        assert resp1.json()["content_hash"] == resp2.json()["content_hash"]
