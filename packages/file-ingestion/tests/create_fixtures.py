"""
Create test fixture files for all supported formats.
Run once: python -m tests.create_fixtures
"""

import csv
import os
import sys
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def create_txt():
    """Create a plain text fixture."""
    path = FIXTURES_DIR / "sample.txt"
    path.write_text("Hello, World!\nThis is a test document.\nLine three.", encoding="utf-8")
    return path


def create_md():
    """Create a markdown fixture."""
    path = FIXTURES_DIR / "sample.md"
    path.write_text(
        "# Heading One\n\nFirst section content.\n\n"
        "## Heading Two\n\nSecond section content with details.\n\n"
        "### Sub-heading\n\nSub-section text.",
        encoding="utf-8",
    )
    return path


def create_csv():
    """Create a CSV fixture."""
    path = FIXTURES_DIR / "sample.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Age", "City"])
        writer.writerow(["Alice", "30", "New York"])
        writer.writerow(["Bob", "25", "London"])
        writer.writerow(["Charlie", "35", "Tokyo"])
    return path


def create_html():
    """Create an HTML fixture."""
    path = FIXTURES_DIR / "sample.html"
    path.write_text(
        """<!DOCTYPE html>
<html>
<head>
    <title>Test Document</title>
    <meta name="author" content="Test Author">
</head>
<body>
    <h1>Main Heading</h1>
    <p>First paragraph with important content.</p>
    <h2>Sub Heading</h2>
    <p>Second paragraph under sub heading.</p>
</body>
</html>""",
        encoding="utf-8",
    )
    return path


def create_docx():
    """Create a DOCX fixture."""
    try:
        from docx import Document
    except ImportError:
        print("SKIP: python-docx not installed, cannot create DOCX fixture")
        return None

    path = FIXTURES_DIR / "sample.docx"
    doc = Document()
    doc.core_properties.title = "Test DOCX"
    doc.core_properties.author = "Test Author"
    doc.add_heading("Document Title", level=1)
    doc.add_paragraph("This is the first paragraph of the test document.")
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Content under section two.")
    doc.save(str(path))
    return path


def create_pptx():
    """Create a PPTX fixture."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("SKIP: python-pptx not installed, cannot create PPTX fixture")
        return None

    path = FIXTURES_DIR / "sample.pptx"
    prs = Presentation()
    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Test Author"

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Presentation Title"
    slide.placeholders[1].text = "Subtitle text here"

    # Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Content Slide"
    slide.placeholders[1].text = "Bullet point one\nBullet point two\nBullet point three"

    prs.save(str(path))
    return path


def create_xlsx():
    """Create an XLSX fixture."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("SKIP: openpyxl not installed, cannot create XLSX fixture")
        return None

    path = FIXTURES_DIR / "sample.xlsx"
    wb = Workbook()
    wb.properties.title = "Test Spreadsheet"
    wb.properties.creator = "Test Author"

    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value", "Category"])
    ws.append(["Alpha", 100, "A"])
    ws.append(["Beta", 200, "B"])
    ws.append(["Gamma", 300, "A"])

    ws2 = wb.create_sheet("Summary")
    ws2.append(["Total", 600])

    wb.save(str(path))
    return path


def create_pdf():
    """Create a PDF fixture using a minimal PDF structure (no extra dependencies)."""
    path = FIXTURES_DIR / "sample.pdf"

    # Minimal valid PDF with text content
    # This creates a simple single-page PDF with embedded text
    content = """%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj

2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj

3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj

4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Hello PDF World) Tj ET
endstream
endobj

5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj

xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n

trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""

    path.write_text(content, encoding="ascii")
    return path


def main():
    """Create all fixture files."""
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)

    creators = [
        ("TXT", create_txt),
        ("MD", create_md),
        ("CSV", create_csv),
        ("HTML", create_html),
        ("DOCX", create_docx),
        ("PPTX", create_pptx),
        ("XLSX", create_xlsx),
        ("PDF", create_pdf),
    ]

    for name, creator in creators:
        result = creator()
        if result:
            print(f"Created {name}: {result} ({result.stat().st_size} bytes)")
        else:
            print(f"SKIP: {name}")

    print(f"\nAll fixtures in: {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
