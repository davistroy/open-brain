#!/usr/bin/env python3
"""
Repair and retry failed file extractions.

Reads error records from the ingestion progress DB, classifies each failure,
and attempts extraction using the appropriate fallback method:

  - .doc files → LibreOffice headless conversion to .docx, then python-docx
  - Large/corrupt PDFs → pymupdf (fitz), then pdftotext (poppler), then OCR
  - Collateral crash files → simple retry via primary extractor
  - Chartsheet XLSX → skip chart sheets, extract data sheets only
  - Encrypted files → log and skip (need password)

Runs in Docker on homeserver with all tools pre-installed.

Usage (inside Docker container):
    python ingest-repair.py --status
    python ingest-repair.py --dry-run
    python ingest-repair.py
"""

from __future__ import annotations

import argparse
import json
import logging
import os
import sqlite3
import subprocess
import sys
import tempfile
import time
from collections import Counter
from pathlib import Path
from typing import Any

import requests

sys.stdout.reconfigure(line_buffering=True)  # type: ignore[attr-defined]
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("ingest-repair")

PROGRESS_DB = os.environ.get("PROGRESS_DB", "/progress/repair-progress.db")
ERRORS_JSON = os.environ.get("ERRORS_JSON", "/progress/repair-errors.json")
ONEDRIVE_ROOT = os.environ.get("ONEDRIVE_ROOT", "/onedrive")
CORE_API = os.environ.get("CORE_API_URL", "http://open-brain-core-api:3000")
CALLER_HEADER = "X-Open-Brain-Caller"
CALLER_VALUE = "file-ingestion-repair"

DOMAIN_TO_VIEW: dict[str, str] = {
    "Work": "work-internal",
    "Amateur Radio": "technical",
    "Making": "technical",
    "Personal": "personal",
    "Sailing": "personal",
    "Projects": "technical",
    "Reference": "technical",
    "App Data": "technical",
    "_Archive": "personal",
}


def classify_error(path: str, error: str | None) -> str:
    """Classify an error into a repair strategy."""
    ext = Path(path).suffix.lower()
    error = error or ""

    if "Unsupported file type: .doc" in error:
        return "doc_convert"
    if "BadZipFile" in error:
        return "xlsx_repair"
    if "CDFV2 Encrypted" in error or "encrypted" in error.lower():
        return "encrypted_skip"
    if "Chartsheet" in error or "iter_rows" in error:
        return "xlsx_chartsheet"
    if "PdfminerException" in error:
        return "pdf_fallback"
    if "extraction timeout" in error or "file too large" in error:
        if ext == ".pdf":
            return "pdf_fallback"
        return "large_fallback"
    if "Connection refused" in error or "Connection aborted" in error or "Max retries" in error:
        return "retry"
    if "Read timed out" in error:
        if ext == ".pdf":
            return "pdf_fallback"
        return "retry"
    return "unknown"


# ── Extractors ──────────────────────────────────────────────────────────────


def extract_doc_via_libreoffice(filepath: str) -> tuple[str | None, str | None]:
    """Convert .doc to .docx via LibreOffice, then extract with python-docx."""
    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, filepath],
            capture_output=True,
            text=True,
            timeout=90,
        )
        if result.returncode != 0:
            return None, f"libreoffice failed: {result.stderr[:200]}"
        docx_files = list(Path(tmpdir).glob("*.docx"))
        if not docx_files:
            return None, "libreoffice produced no output"
        try:
            from docx import Document  # type: ignore[import-untyped]

            doc = Document(str(docx_files[0]))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text, None
        except Exception as e:
            return None, f"python-docx failed: {e}"


def extract_pdf_pymupdf(filepath: str) -> tuple[str | None, str | None]:
    """Extract text from PDF using pymupdf (fast C-based parser)."""
    try:
        import fitz  # type: ignore[import-untyped]

        doc = fitz.open(filepath)
        pages: list[str] = []
        for i, page in enumerate(doc):  # type: ignore[arg-type]
            if i >= 500:  # safety cap
                break
            text: str = page.get_text()
            if text.strip():
                pages.append(text)
        doc.close()
        return "\n\n".join(pages), None
    except Exception as e:
        return None, f"pymupdf failed: {e}"


def extract_pdf_pdftotext(filepath: str) -> tuple[str | None, str | None]:
    """Extract text using poppler's pdftotext (very fast, handles most PDFs)."""
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", filepath, "-"], capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout, None
        return None, f"pdftotext failed: rc={result.returncode}"
    except subprocess.TimeoutExpired:
        return None, "pdftotext timeout"
    except FileNotFoundError:
        return None, "pdftotext not installed"


def extract_pdf_ocr(filepath: str) -> tuple[str | None, str | None]:
    """OCR a scanned PDF using tesseract via pymupdf."""
    try:
        import fitz  # type: ignore[import-untyped]

        doc = fitz.open(filepath)
        pages: list[str] = []
        max_pages = min(len(doc), 100)  # cap OCR at 100 pages
        for i in range(max_pages):
            page = doc[i]
            pix = page.get_pixmap(dpi=200)
            img_path = f"/tmp/ocr_page_{i}.png"
            pix.save(img_path)
            result = subprocess.run(
                ["tesseract", img_path, "stdout", "-l", "eng"],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                pages.append(result.stdout)
            os.unlink(img_path)
        doc.close()
        if pages:
            return "\n\n".join(pages), None
        return None, "OCR produced no text"
    except Exception as e:
        return None, f"OCR failed: {e}"


def extract_xlsx_skip_charts(filepath: str) -> tuple[str | None, str | None]:
    """Extract text from XLSX, skipping chartsheets."""
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
        from openpyxl.chartsheet import Chartsheet  # type: ignore[import-untyped]

        wb = load_workbook(filepath, read_only=True, data_only=True)
        sheets: list[str] = []
        for name in wb.sheetnames:
            ws = wb[name]
            if isinstance(ws, Chartsheet):
                continue
            rows: list[str] = []
            for row in ws.iter_rows(max_row=10000, values_only=True):  # type: ignore[union-attr]
                vals = [str(c) for c in row if c is not None]
                if vals:
                    rows.append("\t".join(vals))
            if rows:
                sheets.append(f"Sheet: {name}\n" + "\n".join(rows))
        wb.close()
        if sheets:
            return "\n\n".join(sheets), None
        return None, "no data sheets found"
    except Exception as e:
        return None, f"xlsx extract failed: {e}"


def extract_xlsx_via_xlrd(filepath: str) -> tuple[str | None, str | None]:
    """Extract text from old .xls binary format using xlrd."""
    try:
        import xlrd  # type: ignore[import-untyped]

        wb = xlrd.open_workbook(filepath)
        sheets: list[str] = []
        for sheet in wb.sheets():
            rows: list[str] = []
            for r in range(min(sheet.nrows, 10000)):
                vals = [
                    str(sheet.cell_value(r, c))
                    for c in range(sheet.ncols)
                    if sheet.cell_value(r, c)
                ]
                if vals:
                    rows.append("\t".join(vals))
            if rows:
                sheets.append(f"Sheet: {sheet.name}\n" + "\n".join(rows))
        if sheets:
            return "\n\n".join(sheets), None
        return None, "no data found"
    except Exception as e:
        return None, f"xlrd failed: {e}"


def extract_large_pptx(filepath: str) -> tuple[str | None, str | None]:
    """Extract text from large PPTX via python-pptx with timeout protection."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(filepath)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides):
            if i >= 200:  # cap
                break
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"Slide {i+1}:\n" + "\n".join(texts))
        if slides:
            return "\n\n".join(slides), None
        return None, "no text found in slides"
    except Exception as e:
        return None, f"pptx extract failed: {e}"


def check_file_magic(filepath: str) -> str:
    """Check file magic bytes to detect format mismatches."""
    try:
        with open(filepath, "rb") as f:
            header = f.read(8)
        if header[:4] == b"\xd0\xcf\x11\xe0":
            return "ole2"  # old binary Office format
        if header[:4] == b"\x50\x4b\x03\x04":
            return "zip"  # modern Office (xlsx/docx/pptx)
        if header[:5] == b"%PDF-":
            return "pdf"
        return "unknown"
    except Exception:
        return "unreadable"


def extract_with_strategy(
    strategy: str, filepath: str, ext: str
) -> tuple[str | None, str | None]:
    """Run the appropriate extraction strategy."""
    if strategy == "doc_convert":
        return extract_doc_via_libreoffice(filepath)

    elif strategy == "pdf_fallback":
        text, err = extract_pdf_pymupdf(filepath)
        if text and len(text.strip()) > 50:
            return text, None
        text, err = extract_pdf_pdftotext(filepath)
        if text and len(text.strip()) > 50:
            return text, None
        # Last resort: OCR
        return extract_pdf_ocr(filepath)

    elif strategy == "xlsx_repair":
        magic = check_file_magic(filepath)
        if magic == "ole2":
            return extract_xlsx_via_xlrd(filepath)
        # Try LibreOffice to repair
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", tmpdir, filepath],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0:
                repaired = list(Path(tmpdir).glob("*.xlsx"))
                if repaired:
                    return extract_xlsx_skip_charts(str(repaired[0]))
        return None, "xlsx repair failed"

    elif strategy == "xlsx_chartsheet":
        return extract_xlsx_skip_charts(filepath)

    elif strategy == "large_fallback":
        if ext in (".pptx",):
            return extract_large_pptx(filepath)
        if ext in (".xlsx", ".xls"):
            return extract_xlsx_skip_charts(filepath)
        # For other large files, try LibreOffice text export
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    tmpdir,
                    filepath,
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode == 0:
                txt_files = list(Path(tmpdir).glob("*.txt"))
                if txt_files:
                    text = txt_files[0].read_text(errors="replace")
                    if text.strip():
                        return text, None
        return None, "large file extraction failed"

    elif strategy == "retry":
        # For collateral crashes, try pymupdf for PDFs, otherwise primary extractor
        if ext == ".pdf":
            text, err = extract_pdf_pymupdf(filepath)
            if text and len(text.strip()) > 50:
                return text, None
            return extract_pdf_pdftotext(filepath)
        elif ext in (".docx",):
            try:
                from docx import Document  # type: ignore[import-untyped]

                doc = Document(filepath)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                return (text, None) if text.strip() else (None, "empty docx")
            except Exception as e:
                return None, f"docx failed: {e}"
        elif ext in (".pptx",):
            return extract_large_pptx(filepath)
        elif ext in (".xlsx",):
            return extract_xlsx_skip_charts(filepath)
        elif ext in (".txt", ".md", ".csv", ".html", ".htm"):
            try:
                text = Path(filepath).read_text(errors="replace")
                return (text, None) if text.strip() else (None, "empty file")
            except Exception as e:
                return None, f"read failed: {e}"
        return None, f"no retry handler for {ext}"

    elif strategy == "encrypted_skip":
        return None, "encrypted — needs password"

    return None, f"unknown strategy: {strategy}"


def submit_capture(
    path: str, text: str, session: requests.Session
) -> tuple[str | None, str | None]:
    """Submit a single file capture to the API."""
    filename = Path(path).name
    top_dir = path.split("/")[0] if "/" in path else "(root)"
    brain_view = DOMAIN_TO_VIEW.get(top_dir, "technical")
    tags: list[str] = [top_dir]
    parts = path.split("/")
    if len(parts) > 1:
        tags.append(parts[1])

    try:
        resp = session.post(
            f"{CORE_API}/api/v1/documents/batch",
            json={
                "files": [
                    {
                        "title": filename,
                        "original_path": path,
                        "mime_type": "application/octet-stream",
                        "brain_view": brain_view,
                        "tags": tags,
                        "content": text[:50000],
                        "category": top_dir,
                        "taxonomy_path": "/".join(parts[:-1]),
                    }
                ]
            },
            timeout=30,
        )
        if resp.status_code in (200, 201):
            data: dict[str, Any] = resp.json()
            results: list[dict[str, Any]] = data.get("results", [])
            cid: str = results[0].get("capture_id", "") if results else ""
            return cid, None
        return None, f"API {resp.status_code}: {resp.text[:200]}"
    except Exception as e:
        return None, str(e)


def init_repair_db() -> sqlite3.Connection:
    conn = sqlite3.connect(PROGRESS_DB, timeout=30)
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("""CREATE TABLE IF NOT EXISTS repair_log (
        path TEXT PRIMARY KEY, status TEXT, error TEXT,
        capture_id TEXT, repaired_at TEXT DEFAULT (datetime('now'))
    )""")
    conn.commit()
    return conn


def run_repair(args: argparse.Namespace) -> None:
    conn = init_repair_db()
    session = requests.Session()
    session.headers[CALLER_HEADER] = CALLER_VALUE

    with open(ERRORS_JSON) as f:
        errors: list[list[str]] = json.load(f)

    # Skip already repaired
    done: set[str] = set(r[0] for r in conn.execute("SELECT path FROM repair_log").fetchall())
    errors = [(p, e) for p, e in errors if p not in done]  # type: ignore[assignment]
    log.info(f"Skipping {len(done)} already attempted, {len(errors)} remaining")

    # Classify errors
    strategies: Counter[str] = Counter()
    classified: list[tuple[str, str, str]] = []
    for path, error in errors:
        strategy = classify_error(path, error)
        strategies[strategy] += 1
        classified.append((path, error, strategy))

    log.info(f"Total errors: {len(errors)}")
    log.info("Strategies:")
    for s, c in strategies.most_common():
        log.info(f"  {s}: {c}")

    if args.dry_run:
        log.info("DRY RUN — no extraction or submission")
        return

    repaired = 0
    still_failed = 0
    skipped = 0

    for i, (path, _orig_error, strategy) in enumerate(classified):
        if strategy == "encrypted_skip":
            skipped += 1
            continue
        if strategy == "unknown":
            still_failed += 1
            continue

        filepath = os.path.join(ONEDRIVE_ROOT, path)
        if not os.path.exists(filepath):
            still_failed += 1
            continue

        ext = Path(path).suffix.lower()
        text, err = extract_with_strategy(strategy, filepath, ext)

        if not text or len(text.strip()) < 50:
            conn.execute(
                "INSERT OR REPLACE INTO repair_log(path, status, error) VALUES(?,?,?)",
                (path, "failed", err or "empty after repair"),
            )
            conn.commit()
            still_failed += 1
            if (i + 1) % 50 == 0:
                log.info(
                    f"  [{i+1}/{len(classified)}] repaired={repaired} failed={still_failed} skipped={skipped}"
                )
            continue

        # Submit to API
        cid, submit_err = submit_capture(path, text, session)
        if cid:
            conn.execute(
                "INSERT OR REPLACE INTO repair_log(path, status, capture_id) VALUES(?,?,?)",
                (path, "ok", cid),
            )
            repaired += 1
        else:
            conn.execute(
                "INSERT OR REPLACE INTO repair_log(path, status, error) VALUES(?,?,?)",
                (path, "submit_failed", submit_err or "submit failed"),
            )
            still_failed += 1

        conn.commit()

        if (i + 1) % 50 == 0:
            log.info(
                f"  [{i+1}/{len(classified)}] repaired={repaired} failed={still_failed} skipped={skipped}"
            )

        time.sleep(0.1)

    log.info("\n=== REPAIR COMPLETE ===")
    log.info(f"  Repaired: {repaired}")
    log.info(f"  Still failed: {still_failed}")
    log.info(f"  Skipped (encrypted): {skipped}")
    conn.close()


def show_status() -> None:
    with open(ERRORS_JSON) as f:
        errors: list[list[str]] = json.load(f)
    strategies: Counter[str] = Counter()
    for path, error in errors:
        strategies[classify_error(path, error)] += 1
    print("\n=== Repair Status ===")
    print(f"Total errors to repair: {len(errors)}")
    for s, c in strategies.most_common():
        print(f"  {s}: {c}")

    conn = init_repair_db()
    repaired: int = conn.execute(
        "SELECT COUNT(*) FROM repair_log WHERE status='ok'"
    ).fetchone()[0]
    failed: int = conn.execute(
        "SELECT COUNT(*) FROM repair_log WHERE status='failed'"
    ).fetchone()[0]
    submit_failed: int = conn.execute(
        "SELECT COUNT(*) FROM repair_log WHERE status='submit_failed'"
    ).fetchone()[0]
    print(f"\nRepaired: {repaired}")
    print(f"Failed: {failed}")
    print(f"Submit failed: {submit_failed}")
    print(f"Remaining: {len(errors) - repaired - failed - submit_failed}")
    print()
    conn.close()


def main() -> None:
    ap = argparse.ArgumentParser(description="Repair failed file extractions")
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--status", action="store_true")
    args = ap.parse_args()

    if args.status:
        show_status()
        return
    run_repair(args)


if __name__ == "__main__":
    main()
